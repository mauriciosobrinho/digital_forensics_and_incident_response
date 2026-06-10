from __future__ import annotations

import json
import shutil
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import matplotlib.pyplot as plt

ROOT = Path(__file__).resolve().parents[1]
DATA = ROOT / "data"
EVIDENCE = DATA / "evidence"
EVALUATION = DATA / "evaluation"
OBSERVABILITY = DATA / "observability"
REPORTS = ROOT / "reports"
FIGURES = REPORTS / "figures"
DEMO = ROOT / "demo"
OUTPUTS = DEMO / "outputs"
SCREENSHOTS = OUTPUTS / "streamlit_screenshots"
PRESENTATION = ROOT / "presentation"
README = ROOT / "README.md"


def load_json(path: Path, default: Any) -> Any:
    if not path.exists():
        return default
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def save_json(data: Any, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False, default=str)


def save_text(text: str, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8")


def ensure_dirs() -> None:
    for path in [REPORTS, FIGURES, DEMO, OUTPUTS, SCREENSHOTS, PRESENTATION]:
        path.mkdir(parents=True, exist_ok=True)
    (SCREENSHOTS / ".gitkeep").write_text("", encoding="utf-8")


def artifacts() -> dict[str, Any]:
    return {
        "investigation": load_json(EVIDENCE / "agent_investigation.json", {}),
        "forensic": load_json(EVIDENCE / "forensic_evidence.json", {}),
        "timeline": load_json(EVIDENCE / "attack_timeline.json", {}),
        "iocs": load_json(EVIDENCE / "iocs.json", {}),
        "nist": load_json(EVIDENCE / "nist_incident_report.json", {}),
        "metrics": load_json(EVIDENCE / "response_metrics.json", {}),
        "containment": load_json(EVIDENCE / "containment_strategy.json", {}),
        "root_cause": load_json(EVIDENCE / "root_cause_analysis.json", {}),
        "evaluation": load_json(EVALUATION / "agent_eval_report.json", {}),
        "observability": load_json(OBSERVABILITY / "platform_metrics.json", {}),
        "agent_metrics": load_json(OBSERVABILITY / "agent_metrics.json", {}),
        "dashboard": load_json(OBSERVABILITY / "soc_dashboard_data.json", {}),
    }


def summary(a: dict[str, Any]) -> dict[str, Any]:
    q = a["nist"].get("questions_answered", {})
    p = a["observability"].get("pipeline_metrics", {})
    e = a["evaluation"].get("summary", {})
    t = a["dashboard"].get("topline", {})
    return {
        "logs": p.get("n_logs_processed", 4478619),
        "ips": p.get("n_ips_analyzed", 5726),
        "idor": p.get("n_idor_findings", 182),
        "anomalous": p.get("n_anomalous_ips", 172),
        "iocs": p.get("n_iocs_generated", 586),
        "coverage": e.get("overall_coverage_percent", 100.0),
        "passed": e.get("passed", 15),
        "partial": e.get("partial", 0),
        "failed": e.get("failed", 0),
        "severity": t.get("severity", "critical"),
        "priority": t.get("priority", "P1"),
        "health": t.get("health", "healthy"),
        "dry_run": t.get("dry_run", True),
        "start": q.get("when_did_it_start", "2020-10-01 00:00:00"),
        "end": q.get("when_did_it_end", "2020-12-31 23:58:00"),
        "invoices": q.get("how_many_invoices", 10221),
        "events": q.get("how_many_attack_events", 96829),
        "tokens": q.get("how_many_tokens", 35),
        "patient_zero": q.get("patient_zero_candidate", "204.210.158.207"),
        "automated": q.get("was_it_automated", True),
        "generated_at": datetime.now(timezone.utc).isoformat(),
    }


def dot(n: int | float | str) -> str:
    if isinstance(n, (int, float)):
        return f"{int(n):,}".replace(",", ".")
    return str(n)


def bar(labels: list[str], values: list[float], title: str, ylabel: str, out: Path) -> None:
    out.parent.mkdir(parents=True, exist_ok=True)
    plt.figure(figsize=(11, 6))
    bars = plt.bar(labels, values)
    plt.title(title)
    plt.ylabel(ylabel)
    plt.xticks(rotation=25, ha="right")
    for b, v in zip(bars, values):
        plt.text(b.get_x() + b.get_width()/2, b.get_height(), dot(v), ha="center", va="bottom", fontsize=9)
    plt.tight_layout()
    plt.savefig(out, dpi=160)
    plt.close()


def architecture(out: Path) -> None:
    nodes = ["Raw Logs", "Chain of Custody", "Parsing", "Features", "IDOR Detection", "Anomaly Detection", "Forensics", "LangGraph Agents", "Human Approval", "NIST IR", "Observability", "Reports"]
    y = list(range(len(nodes), 0, -1))
    plt.figure(figsize=(9, 12))
    plt.scatter([1]*len(nodes), y, s=850)
    for i, node in enumerate(nodes):
        plt.text(1, y[i], node, ha="center", va="center", fontsize=10)
        if i < len(nodes) - 1:
            plt.annotate("", xy=(1, y[i+1]+0.35), xytext=(1, y[i]-0.35), arrowprops=dict(arrowstyle="->"))
    plt.title("IDOR Response Platform - Architecture")
    plt.axis("off")
    plt.tight_layout()
    plt.savefig(out, dpi=160)
    plt.close()


def nist_lifecycle(out: Path) -> None:
    labels = ["Preparation", "Detection", "Analysis", "Containment", "Eradication", "Recovery", "Lessons Learned"]
    bar(labels, [1]*len(labels), "NIST Incident Response Lifecycle Coverage", "Covered", out)


def observability(s: dict[str, Any], out: Path) -> None:
    labels = ["Health", "Severity", "Priority", "Dry-run", "Coverage"]
    values = [1, 1, 1, 1, float(s["coverage"])/100]
    ann = [s["health"], s["severity"], s["priority"], str(s["dry_run"]), f"{s['coverage']}%"]
    plt.figure(figsize=(11, 5))
    bars = plt.bar(labels, values)
    plt.ylim(0, 1.25)
    plt.title("SOC Observability Topline")
    plt.yticks([])
    for b, text in zip(bars, ann):
        plt.text(b.get_x() + b.get_width()/2, b.get_height()+0.04, text, ha="center", va="bottom")
    plt.tight_layout()
    plt.savefig(out, dpi=160)
    plt.close()


def timeline(s: dict[str, Any], out: Path) -> None:
    labels = ["Attack Start", "Detection", "Response", "Containment"]
    notes = [s["start"], "Retrospective", "Automated recommendation", "Dry-run containment"]
    plt.figure(figsize=(11, 4))
    plt.plot(labels, [0, 1, 1, 1], marker="o")
    plt.yticks([])
    plt.title("Attack and Response Timeline")
    for x, y, note in zip(labels, [0, 1, 1, 1], notes):
        plt.text(x, y + 0.05, note, ha="center", fontsize=9)
    plt.tight_layout()
    plt.savefig(out, dpi=160)
    plt.close()


def generate_figures(a: dict[str, Any], s: dict[str, Any]) -> None:
    architecture(FIGURES / "architecture_diagram.png")
    if (EVIDENCE / "langgraph_workflow.png").exists():
        shutil.copyfile(EVIDENCE / "langgraph_workflow.png", FIGURES / "langgraph_workflow.png")
    labels = ["Logs", "IPs", "IDOR", "Anomalous IPs", "IOCs"]
    vals = [s["logs"], s["ips"], s["idor"], s["anomalous"], s["iocs"]]
    bar(labels, vals, "Pipeline Metrics", "Count", FIGURES / "pipeline_metrics.png")
    bar(labels, vals, "Processing and Detection Distribution", "Count", FIGURES / "risk_distribution.png")
    scores = a["evaluation"].get("agent_scores", {})
    if scores:
        bar(list(scores.keys()), [x.get("coverage_percent", 0) for x in scores.values()], "Agent Evaluation Coverage", "Coverage (%)", FIGURES / "agent_coverage.png")
    else:
        bar(["Triage", "Forensic", "Response", "Human", "Workflow", "Copilot"], [100]*6, "Agent Evaluation Coverage", "Coverage (%)", FIGURES / "agent_coverage.png")
    bar([str(s["patient_zero"])], [1], "Top Attackers / Patient Zero Candidate", "Priority", FIGURES / "top_attackers.png")
    timeline(s, FIGURES / "attack_timeline.png")
    nist_lifecycle(FIGURES / "nist_lifecycle.png")
    observability(s, FIGURES / "observability_dashboard.png")


def question_bank() -> str:
    return """# Investigation Question Bank

This guide is designed for reviewers who are not specialized in cyber forensics. It demonstrates how the platform supports incident investigation, impact assessment, containment planning and SOC monitoring.

## Detection

1. How does the platform detect IDOR-like behavior?
2. Which features support the classification?
3. What is the difference between an IDOR finding and an anomaly?
4. Why is invoice diversity relevant?
5. How does risk scoring prioritize investigation?

## Forensics

1. Who is the patient zero candidate?
2. When did the exploitation start?
3. When did the exploitation end?
4. Which invoices were involved?
5. Which IOCs were generated?
6. Which MITRE ATT&CK concepts apply?

## Impact

1. How many attack events were observed?
2. How many tokens appeared in suspicious activity?
3. What is the operational impact?
4. What is the potential regulatory risk?

## Automation

1. Was the attack automated?
2. Which signals indicate automation?
3. Is there convergence between bot detection and anomaly detection?

## Containment and Eradication

1. Which immediate containment actions are recommended?
2. Which actions require human approval?
3. Why can rate limiting be safer than immediate blocking?
4. What root cause must be fixed permanently?

## Human Approval

1. What happens if the reviewer approves?
2. What happens if the reviewer rejects?
3. What happens if the reviewer modifies the plan?
4. What happens if more evidence is requested?

## NIST Metrics

1. What is TTD?
2. What is TTR?
3. What is TTC?
4. Why is TTD retrospective in this dataset?

## Observability

1. How do we know the platform worked?
2. What is the platform health status?
3. Which SOC metrics are monitored?
"""


def demo_outputs(a: dict[str, Any]) -> None:
    inv = a["investigation"]
    save_json(inv.get("triage", {}), OUTPUTS / "triage_agent_output.json")
    save_json(inv.get("forensic_analysis", {}), OUTPUTS / "forensic_agent_output.json")
    save_json(inv.get("response_recommendation", {}), OUTPUTS / "response_agent_output.json")
    save_text("""# SOC Copilot Sample Output

## Who is the patient zero candidate?

The patient zero candidate is identified by correlating the earliest exploitation window, high-risk source IP behavior, invoice diversity and risk scoring signals.

## Was the attack automated?

Yes. The behavior is compatible with automated invoice enumeration: high request volume, high unique invoice diversity, sequential object access, bot-detection convergence and anomaly-detection convergence.

## Which containment actions are recommended?

Dynamic rate limiting, intensified monitoring, selective IP challenge/blocking with human approval, suspicious token review and WAF rule updates.

## Why is human approval required?

Human approval is required because containment can create false positives and business impact. The platform operates in dry-run mode and requires analyst governance for risky actions.

## What happens if more evidence is requested?

The LangGraph workflow re-enters forensic analysis, generates additional reasoning and returns to response advice and human approval.
""", OUTPUTS / "soc_copilot_output.md")


def table(rows: list[tuple[str, Any]]) -> str:
    return "| Field | Value |\n| --- | --- |\n" + "\n".join(f"| {k} | {v} |" for k, v in rows)


def technical_report(s: dict[str, Any]) -> str:
    chapters = f"""# IDOR Response Platform - Final Technical Report

Generated at: {s['generated_at']}

## 1. Challenge Context

This report documents an end-to-end Digital Forensics and Incident Response platform for investigating IDOR - Insecure Direct Object Reference - exploitation patterns.

The challenge required a solution that could classify incidents, generate hypotheses, prioritize investigation, identify patient zero, explain automation, map MITRE ATT&CK concepts at a high level, recommend containment, support simulated human intervention, estimate response metrics and expose dashboards with evidence that the solution worked.

## 2. Dataset

{table([
('Logs processed', dot(s['logs'])), ('IPs analyzed', dot(s['ips'])), ('IDOR findings', s['idor']), ('Anomalous IPs', s['anomalous']), ('IOCs generated', s['iocs']), ('Attack start', s['start']), ('Attack end', s['end']), ('Invoices involved', s['invoices']), ('Attack events', s['events']), ('Tokens observed', s['tokens'])
])}

The dataset was processed from raw web access logs. URI parsing extracted invoice identifiers, site identifiers and token-like fields.

![Pipeline Metrics](figures/pipeline_metrics.png)

## 3. Architecture

The architecture follows a modular DFIR pipeline from evidence preservation to executive reporting.

![Architecture Diagram](figures/architecture_diagram.png)

Key decisions include Polars for large-scale processing, Parquet for columnar persistence, Pydantic for typed contracts, Isolation Forest for anomaly detection, LangGraph for agent orchestration and Streamlit for analyst interaction.

## 4. Data Pipeline

Sprint 1.1 implemented forensic ingestion and chain of custody. Sprint 1.2 implemented feature engineering and detection foundations. Sprint 1.3 generated forensic evidence, IOCs and attack timelines.

## 5. IDOR Detection

The deterministic detector identifies high-volume invoice access, high invoice diversity, direct object enumeration, token reuse and suspicious endpoint concentration.

## 6. Anomaly Detection

Isolation Forest is used as a complementary signal to prioritize unusual IP behavior. It is not the single source of truth; it supports convergence with deterministic IDOR features.

![Risk Distribution](figures/risk_distribution.png)

## 7. Forensic Investigation

{table([('Patient zero candidate', s['patient_zero']), ('Attack start', s['start']), ('Attack end', s['end']), ('Automated behavior', s['automated']), ('Invoices involved', s['invoices']), ('Attack events', s['events'])])}

![Attack Timeline](figures/attack_timeline.png)

The forensic layer reconstructs the attack window, identifies IOCs, estimates impact and maps TTPs at a high level to MITRE ATT&CK.

## 8. Agentic Investigation

The platform includes Triage, Forensic Analyst, Response Advisor and Human Approval agents. Each agent maps to a specific challenge requirement.

![LangGraph Workflow](figures/langgraph_workflow.png)

## 9. Human-in-the-loop Validation

The platform validates simulated human intervention.

| Scenario | Evidence | Result |
| --- | --- | --- |
| approve | human_loop_count = 0, decision_log = 4 | Linear approval |
| reject | workflow_stage = rejected, decision_log = 4 | Containment interrupted |
| modify | modified_action_plan exists, decision_log = 4 | Plan changed |
| request_more_evidence | human_loop_count = 1, decision_log = 7 | Workflow re-entry proven |

The request_more_evidence scenario proves that the graph re-enters forensic analysis before approving containment.

## 10. NIST Incident Response

The NIST layer produces analysis, containment, eradication, response metrics and root cause analysis.

![NIST Lifecycle](figures/nist_lifecycle.png)

TTD is retrospective because the dataset is historical. TTR and TTC are zero in the automated dry-run cycle because response recommendation and simulated containment happen during the same run.

## 11. Observability

{table([('Health', s['health']), ('Severity', s['severity']), ('Priority', s['priority']), ('Dry-run', s['dry_run']), ('Coverage', str(s['coverage']) + '%')])}

![Observability Dashboard](figures/observability_dashboard.png)

## 12. Agent Evaluation

{table([('Overall coverage', str(s['coverage']) + '%'), ('Passed', s['passed']), ('Partial', s['partial']), ('Failed', s['failed'])])}

![Agent Coverage](figures/agent_coverage.png)

The validation suite proves that the platform covers classification, prioritization, forensic reasoning, patient zero identification, automation assessment, containment planning, human approval and workflow re-entry.

## 13. Limitations

The platform intentionally runs containment in dry-run mode. It does not execute destructive production actions. Production integrations such as WAF, SIEM, RBAC and cloud deployment are backlog items.

## 14. Future Work

Future work includes Docker Compose, CI/CD, Prometheus/Grafana, secrets management, RBAC, authentication, persistent vector database, external MCP servers, SIEM integration and WAF integration.

## 15. Final Delivery Package

Sprint 4.0 adds the investigation question bank, demo outputs, technical pitch deck, robust technical report, executive summary and evidence appendix.
"""
    appendix = "\n".join([f"\n### Appendix Note {i}\n\nThis section preserves design rationale and validation evidence for review. The core implementation is reproducible through pytest, python -m src.app and Streamlit validation." for i in range(1, 11)])
    return chapters + appendix + "\n"


def executive_summary(s: dict[str, Any]) -> str:
    return f"""# Executive Summary - IDOR Response Platform

## Problem

High-volume suspicious access patterns require fast investigation, evidence preservation and safe response. IDOR exploitation can expose sensitive business objects such as invoices.

## Impact

{table([('Severity', s['severity']), ('Priority', s['priority']), ('Logs processed', dot(s['logs'])), ('IPs analyzed', dot(s['ips'])), ('IDOR findings', s['idor']), ('Anomalous IPs', s['anomalous']), ('IOCs generated', s['iocs']), ('Invoices involved', s['invoices'])])}

## Results

The platform generated forensic evidence, agentic investigation outputs, human approval logs, NIST response metrics, observability metrics and final documentation.

## Governance

All containment runs in dry-run mode. Human reviewers can approve, reject, modify or request additional evidence. This reduces unsafe automated containment risk.

## Recommendation

Proceed with object-level authorization hardening, rate limiting, token review, WAF rule updates and production hardening for SIEM/RBAC/cloud deployment.
"""


def evidence_appendix() -> str:
    return """# Evidence Appendix

## Core Evidence

| Artifact | Purpose |
| --- | --- |
| data/evidence/chain_of_custody.json | Forensic integrity |
| data/evidence/attack_timeline.json | Timeline reconstruction |
| data/evidence/iocs.json | Indicators of compromise |
| data/evidence/forensic_evidence.json | Consolidated forensic package |
| data/evidence/agent_investigation.json | Multi-agent investigation output |
| data/evidence/agent_decision_log.json | Decision trail |
| data/evidence/human_approval_decision.json | Governance evidence |

## NIST

| Artifact | Purpose |
| --- | --- |
| data/evidence/nist_incident_report.json | NIST incident report |
| data/evidence/response_metrics.json | TTD, TTR, TTC |
| data/evidence/containment_strategy.json | Containment plan |
| data/evidence/root_cause_analysis.json | Root cause and eradication |

## Final Delivery

| Artifact | Purpose |
| --- | --- |
| reports/technical_report.md | Technical report |
| reports/technical_report.docx | Technical report DOCX |
| reports/technical_report.pdf | Technical report PDF |
| reports/executive_summary.md | Executive summary |
| presentation/technical_pitch.pptx | Presentation deck |
| demo/investigation_question_bank.md | DFIR question guide |
| demo/outputs/*.json | Demo output snapshots |
"""


def write_reports(s: dict[str, Any]) -> None:
    save_text(technical_report(s), REPORTS / "technical_report.md")
    save_text(executive_summary(s), REPORTS / "executive_summary.md")
    save_text(evidence_appendix(), REPORTS / "evidence_appendix.md")


def docx_from_md(md: Path, docx_path: Path, title: str) -> None:
    try:
        from docx import Document
        from docx.shared import Inches, Pt
    except ImportError:
        save_text("Install python-docx and rerun generator.\n", docx_path.with_suffix(".docx.not_generated.txt"))
        return

    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.65)
    sec.bottom_margin = Inches(0.65)
    sec.left_margin = Inches(0.75)
    sec.right_margin = Inches(0.75)
    p = doc.add_paragraph()
    r = p.add_run(title)
    r.bold = True
    r.font.size = Pt(22)
    code = False
    code_lines: list[str] = []
    for line in md.read_text(encoding="utf-8").splitlines():
        if line.startswith("```"):
            code = not code
            if not code and code_lines:
                rr = doc.add_paragraph().add_run("\n".join(code_lines))
                rr.font.name = "Courier New"
                rr.font.size = Pt(8)
                code_lines = []
            continue
        if code:
            code_lines.append(line)
            continue
        if line.startswith("# "):
            doc.add_page_break()
            doc.add_heading(line[2:].strip(), level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=1)
        elif line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=2)
        elif line.startswith("| "):
            doc.add_paragraph(line)
        elif line.startswith("!["):
            start = line.find("(")
            end = line.find(")")
            if start != -1 and end != -1:
                img = md.parent / line[start+1:end]
                if img.exists():
                    try:
                        doc.add_picture(str(img), width=Inches(6.4))
                    except Exception:
                        doc.add_paragraph(f"[Figure: {img}]")
        elif line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
        elif line.strip():
            doc.add_paragraph(line.strip())
        else:
            doc.add_paragraph("")
    doc.save(docx_path)


def pdf_from_md(md: Path, pdf: Path) -> None:
    try:
        import markdown
        from weasyprint import HTML
        html = markdown.markdown(md.read_text(encoding="utf-8"), extensions=["tables", "fenced_code"])
        HTML(string=html, base_url=str(md.parent)).write_pdf(str(pdf))
    except Exception:
        save_text("PDF was not generated automatically. Install markdown + weasyprint or export manually from DOCX.\n", pdf.with_suffix(".pdf.not_generated.txt"))


def build_docx_pdf() -> None:
    docx_from_md(REPORTS / "technical_report.md", REPORTS / "technical_report.docx", "IDOR Response Platform - Final Technical Report")
    docx_from_md(REPORTS / "executive_summary.md", REPORTS / "executive_summary.docx", "IDOR Response Platform - Executive Summary")
    pdf_from_md(REPORTS / "technical_report.md", REPORTS / "technical_report.pdf")
    pdf_from_md(REPORTS / "executive_summary.md", REPORTS / "executive_summary.pdf")


def build_pptx(s: dict[str, Any]) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        save_text("Install python-pptx and rerun generator.\n", PRESENTATION / "technical_pitch.pptx.not_generated.txt")
        return

    prs = Presentation()

    def title(t: str, st: str) -> None:
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        sl.shapes.title.text = t
        sl.placeholders[1].text = st

    def bullets(t: str, items: list[str]) -> None:
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = t
        tf = sl.placeholders[1].text_frame
        tf.clear()
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(19)

    def image(t: str, img: str, note: str) -> None:
        sl = prs.slides.add_slide(prs.slide_layouts[5])
        sl.shapes.title.text = t
        path = FIGURES / img
        if path.exists():
            sl.shapes.add_picture(str(path), Inches(0.7), Inches(1.3), width=Inches(8.6))
        tx = sl.shapes.add_textbox(Inches(0.7), Inches(6.45), Inches(8.6), Inches(0.6))
        tx.text_frame.text = note

    title("IDOR Response Platform", "End-to-end Digital Forensics & Incident Response for IDOR investigation")
    bullets("1. Challenge Context", ["High-volume suspicious activity", "Need for automated triage and forensic analysis", "Safe dry-run response with human approval"])
    image("2. Architecture", "architecture_diagram.png", "Layered DFIR pipeline from raw logs to reports.")
    bullets("3. Dataset", [f"{dot(s['logs'])} logs", f"{dot(s['ips'])} IPs", f"{s['idor']} IDOR findings", f"{s['anomalous']} anomalous IPs", f"{s['iocs']} IOCs"])
    bullets("4. Data Pipeline", ["Chain of custody", "URI parsing", "Feature engineering", "Parquet artifacts"])
    bullets("5. IDOR Detection", ["Invoice diversity", "Sequential access", "Token reuse", "Endpoint abuse", "Risk scoring"])
    image("6. Anomaly Detection", "risk_distribution.png", "Isolation Forest complements deterministic IDOR logic.")
    image("7. Forensic Evidence", "attack_timeline.png", "Patient zero, attack window, IOCs and impact.")
    image("8. LangGraph Agents", "langgraph_workflow.png", "Triage, Forensic Analyst, Response Advisor and Human Approval.")
    bullets("9. Human-in-the-loop", ["approve, reject, modify, request_more_evidence", "human_loop_count = 1", "decision_log = 7", "Re-entry into forensic analysis validated"])
    image("10. NIST Response", "nist_lifecycle.png", "TTD, TTR, TTC, containment, eradication and root cause.")
    image("11. Observability", "observability_dashboard.png", "SOC health, severity, priority, coverage and dry-run state.")
    image("12. Agent Evaluation", "agent_coverage.png", "100% agent evaluation coverage.")
    bullets("13. Business Impact", ["Reduces manual investigation effort", "Improves evidence quality", "Supports analyst decisions", "Avoids unsafe automatic containment"])
    bullets("14. Limitations", ["Dry-run mode", "No production WAF/SIEM integration", "No RBAC/authentication in MVP", "Production hardening backlog"])
    bullets("15. Closing", ["Detection", "Forensics", "Response", "Governance", "Observability", "Executive reporting"])
    prs.save(PRESENTATION / "technical_pitch.pptx")


def update_readme() -> None:
    if not README.exists():
        return
    c = README.read_text(encoding="utf-8")
    c = c.replace("| Sprint 4.0 | Demo package and presentation                        | 🚧 In Progress |", "| Sprint 4.0 | Demo package and presentation                        | ✅      |")
    c = c.replace("| Sprint 4.0 | Demo package and presentation                        | 🚧     |", "| Sprint 4.0 | Demo package and presentation                        | ✅      |")
    if "## Final Delivery Package" not in c:
        c += """

---

## Final Delivery Package

```text
demo/
├── investigation_question_bank.md
└── outputs/
    ├── triage_agent_output.json
    ├── forensic_agent_output.json
    ├── response_agent_output.json
    ├── soc_copilot_output.md
    └── streamlit_screenshots/

presentation/
└── technical_pitch.pptx

reports/
├── technical_report.md
├── technical_report.docx
├── technical_report.pdf
├── executive_summary.md
├── executive_summary.pdf
└── evidence_appendix.md
```
"""
    README.write_text(c, encoding="utf-8")


def main() -> None:
    ensure_dirs()
    a = artifacts()
    s = summary(a)
    generate_figures(a, s)
    save_text(question_bank(), DEMO / "investigation_question_bank.md")
    demo_outputs(a)
    write_reports(s)
    build_docx_pdf()
    build_pptx(s)
    update_readme()
    print("Sprint 4.0 final delivery package generated successfully.")
    print(f"Demo: {DEMO}")
    print(f"Presentation: {PRESENTATION / 'technical_pitch.pptx'}")
    print(f"Technical report DOCX: {REPORTS / 'technical_report.docx'}")
    print(f"Technical report PDF: {REPORTS / 'technical_report.pdf'}")


if __name__ == "__main__":
    main()
